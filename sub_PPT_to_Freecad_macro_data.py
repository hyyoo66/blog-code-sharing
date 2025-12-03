'''
이름: Large Orange Diamond
코드: 🔶
유니코드: U+1F536
용도: 변경된 부분, 주의점, 하이라이트 표시 등에 사용.
'''
import time
import os
import sys
import re
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_AUTO_SHAPE_TYPE
from pptx.enum.dml import MSO_FILL
from pptx.dml.color import RGBColor, MSO_THEME_COLOR
import logging

logger = logging.getLogger(__name__)

# 로깅 설정
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    handlers=[
#        logging.FileHandler("ppt_processor2.log", encoding='utf-8'),
        logging.StreamHandler(sys.stdout)
    ]
)

logger = logging.getLogger(__name__)
def extract_scale(slide):
    """
    슬라이드에서 @freecad 텍스트 상자를 찾아 scale 값을 추출합니다.
    """
    default_scale = 1.0  # 기본값
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.lower()
            if "@freecad" in text:
                match = re.search(r'scale\s*=\s*([\d\.]+)', text)
                if match:
                    try:
                        scale = float(match.group(1))
                        logger.info(f"scale 값 추출됨: {scale}")
                        return scale
                    except ValueError:
                        logger.warning(f"유효하지 않은 scale 값: {match.group(1)}")
    logger.info("scale 값이 슬라이드에 없으므로 기본값 1.0 사용")
    return default_scale
    

color_map = {
#    '(0:176:80)':'(135:206:235)',  # 비표준 파랑
    '(70:177:225)':'(0:76:80)',   #녹색
    '(233:113:50)':'(255:192:0)',  # 주황색
    '(216:110:204)':'(128:0:128)'  # 보라


}

shape_indices = {
    'RECTANGLE': 9,
    'CIRCLE': 7
}

def adjust_color(fields):
    if len(fields) > 3:
        shape = fields[3]
        if shape in shape_indices:
            index = shape_indices[shape]
            if len(fields) > index and fields[index] in color_map:
                fields[index] = color_map[fields[index]]
    return fields
    
    
def modify_data(line, scale):

    # 공백 제거 및 불필요한 데이터 정리
    line = line.replace(" ", "")
    line = line.replace("ALL,", "")
    line = line.replace("-0.0", "0")
    line = line.replace("0.0", "0")

    # 필드 분리
    fields = line.split(',')

    # z_size가 음수일 경우 처리
    try:
        z_start = float(fields[1])  # z_start 필드
        z_size = float(fields[2])   # z_size 필드
        if z_size == 0:
            fields[2] = str(0.001)
            print(f"z_size가 음수여서 수정됨: z0={fields[1]}, z_size={fields[2]}")
        elif z_size < 0:
            fields[1] = str(z_start + z_size)  # z0를 수정
            fields[2] = str(abs(z_size))
            print(f"z_size가 음수여서 수정됨: z0={fields[1]}, z_size={fields[2]}")
          
    except (ValueError, IndexError) as e:
        print(f"z_size 처리 중 오류 발생: {e}")
        
   
    # 필드 재배열
    if len(fields) >= 11 and fields[4] == 'RECTANGLE':
        field_to_move = fields.pop(3)
        fields.append(field_to_move)
    elif len(fields) >= 9 and fields[4] == 'CIRCLE':
        field_to_move = fields.pop(3)
        fields.append(field_to_move)      

       
    # 색 보정
    fields = adjust_color(fields)

  
    # N 바디에서 11번째 필드를 제거
    if (line.startswith("P") or line.startswith("N")) and len(fields) >= 11 and fields[3].upper() == 'RECTANGLE':
        try:
            del fields[10]  # 11번째 필드 삭제 (0-based index)
            print("11번째 필드를 삭제했습니다.")
        except IndexError:
            print("11번째 필드가 존재하지 않아 삭제할 수 없습니다.")

    if (line.startswith("P") or line.startswith("N")) and len(fields) >= 9 and fields[3].upper() == 'CIRCLE':
        try:
            del fields[8]  # 9번째 필드 삭제 (0-based index)
            print("11번째 필드를 삭제했습니다.")
        except IndexError:
            print("11번째 필드가 존재하지 않아 삭제할 수 없습니다.")

        
    # RECTANGLE 및 CIRCLE 필드에 scale 적용
    try:
        if fields[3] == 'RECTANGLE':
            for i in range(4, 8):  # 4~7번째 필드에 scale 곱하기
                fields[i] = str(round(float(fields[i]) * scale, 3))
        elif fields[3] == 'CIRCLE':
            for i in range(4, 7):  # 4~5번째 필드에 scale 곱하기
                fields[i] = str(round(float(fields[i]) * scale, 3))
    except (ValueError, IndexError) as e:
        logger.warning(f"scale 적용 중 오류 발생: {e}")
 

    # 쉼표를 탭으로 변환
    line = ','.join(fields)
    line = line.replace(",", "\t")

    return line


# EMU 단위를 mm로 변환
def ppt_to_mm(emu_value):
    return round(emu_value / 36000, 3)

# 도형의 중심 좌표 계산
def calculate_center_coordinates(shape, slide_height, x_min, y_min):
    center_x = shape.left + shape.width / 2
    center_y = shape.top + shape.height / 2

    adjusted_x = center_x - x_min
    adjusted_y = slide_height - center_y - y_min  # Y축 반전

    converted_x = ppt_to_mm(adjusted_x)
    converted_y = ppt_to_mm(adjusted_y)
    return converted_x, converted_y

# 슬라이드의 모든 도형에서 최소 x, y 좌표 찾기
def find_min_coordinates(shapes, slide_height):
    min_x = float('inf')
    min_y = float('inf')
    for shape in shapes:
        if shape.shape_type in [MSO_SHAPE_TYPE.AUTO_SHAPE, MSO_SHAPE_TYPE.FREEFORM]:
            min_x = min(min_x, shape.left)
            min_y = min(min_y, slide_height - shape.top - shape.height)  # Y축 반전
    return min_x, min_y  # Correctly return min_y instead of y_min


# RECTANGLE 도형의 회전 각도 계산
def get_shape_rotation(shape):
    return shape.rotation if shape.rotation else 0  # 회전 각도 반환 (기본값 0)

# 실선 여부 확인
def is_solid_line(shape):
    """도형이 실선인지 확인"""
    if not hasattr(shape, 'line') or not shape.line:
        return False  # 도형에 선이 없으면 실선이 아님
    return shape.line.dash_style is None

# 슬라이드에서 z_base 추출
def extract_z_base(slide):
    z_base = 0  # 기본값 설정
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.lower()
            if "@freecad" in text:
                # z_base 값 추출 (예: z_base=10)
                match = re.search(r'z_base\s*=\s*(-?\d+\.?\d*)', text)
                if match:
                    try:
                        z_base = float(match.group(1))
                        logger.info(f"z_base 추출됨: {z_base}")  # z_base 값 로깅
                        return z_base
                    except ValueError:
                        logger.warning(f"유효하지 않은 z_base 값: {match.group(1)}")
    logger.info("z_base 값이 슬라이드에 없으므로 기본값 0 사용")
    return z_base


def get_shape_color(shape):
    """
    PPT 도형에서 RGB 색상 정보를 추출합니다.
    Args:
        shape: PPT 도형 객체
    Returns:
        str: "(R:G:B)" 형식의 색상 문자열
    """
    DEFAULT_COLOR = "(128:128:128)"  # 기본 색상: 회색

    def adjust_brightness(base_rgb, brightness):
        """
        밝기를 기반으로 RGB 값을 조정합니다.
        Args:
            base_rgb: (R, G, B) 튜플
            brightness: 밝기 값 (-1.0 ~ 1.0)
        Returns:
            str: "(R:G:B)" 형식의 조정된 RGB 값
        """
        adjusted = [
            max(0, min(255, int(channel * (1 + brightness))))
            for channel in base_rgb
        ]
        return f"({adjusted[0]}:{adjusted[1]}:{adjusted[2]})"

    try:
        logger.debug("=== Color Extraction Debug ===")
        logger.debug(f"Shape type: {type(shape)}")
        
        # Fill 속성 확인
        if not hasattr(shape, 'fill'):
            logger.debug("Shape has no fill attribute. Using default color.")
            return DEFAULT_COLOR
        
        if shape.fill is None:
            logger.debug("Shape fill is None. Using default color.")
            return DEFAULT_COLOR
        
        # Fill 타입 확인
        logger.debug(f"Fill type: {shape.fill.type}")
        
        if shape.fill.type == MSO_FILL.SOLID:
            logger.debug("Found solid fill")
            fore_color = shape.fill.fore_color
            logger.debug(f"Fore color type: {fore_color.type}")
            
            # RGB 값이 있는 경우
            if hasattr(fore_color, 'rgb') and fore_color.rgb:
                rgb = fore_color.rgb
                logger.debug(f"RGB values: {rgb}")
                if isinstance(rgb, tuple) and len(rgb) == 3:
                    if hasattr(fore_color, 'brightness'):
                        logger.debug(f"Brightness: {fore_color.brightness}")
                        return adjust_brightness(rgb, fore_color.brightness)
                    return f"({rgb[0]}:{rgb[1]}:{rgb[2]})"
            
            # SchemeColor 매핑
            if hasattr(fore_color, 'theme_color'):
                logger.debug(f"Theme color: {fore_color.theme_color}")
                theme_colors = {
                    1: "(0:0:0)",      # TEXT1
                    2: "(255:255:255)", # TEXT2
                    3: "(68:84:106)",   # BACKGROUND1
                    4: "(255:255:255)", # BACKGROUND2
                    5: "(0:176:80)",    # ACCENT1
                    6: "(255:192:0)",   # ACCENT2
                    7: "(255:0:0)",     # ACCENT3
                    8: "(0:112:192)",   # ACCENT4
                    9: "(112:48:160)",  # ACCENT5
                    10: "(0:32:96)",    # ACCENT6
                    14: "(68:68:68)"    # BACKGROUND1 다크 테마
                }
                theme_color = theme_colors.get(fore_color.theme_color, DEFAULT_COLOR)  # 기본값: 회색
                logger.debug(f"Mapped theme color to RGB: {theme_color}")
                return theme_color

            # 밝기(Brightness)가 있는 경우
            if hasattr(fore_color, 'brightness'):
                brightness = fore_color.brightness
                logger.debug(f"Brightness: {brightness}")
                base_rgb = (128, 128, 128)  # 중간 회색 기준으로 조정
                return adjust_brightness(base_rgb, brightness)
        
        # 배경 색상 처리
        elif shape.fill.type == MSO_FILL.BACKGROUND:
            logger.debug("Found background fill. Using default color.")
            return DEFAULT_COLOR

        logger.debug("Using default color (gray)")
        return DEFAULT_COLOR
        
    except Exception as e:
        logger.error(f"Error in color extraction: {str(e)}")
        return DEFAULT_COLOR

def validate_and_adjust_z_property(z_property_original, z_base):
    """
    z_property를 검증하고 z_base 값을 적용하여 수정.
    
    :param z_property_original: 원본 z_property 문자열
    :param z_base: 기준값
    :return: 수정된 z_property 문자열 또는 None (유효하지 않을 경우)
    """
    z_prop_parts = [part.strip() for part in z_property_original.split(",")]

    # 필드 개수 검사 (최소 2개 이상, 최대 4개)
    if len(z_prop_parts) < 2 or len(z_prop_parts) > 4:
        logger.warning(f"z_property 필드 개수가 유효하지 않음: {z_property_original}")
        return None

    # 0번 필드 검사 (d, p, n 중 하나, 대소문자 무시)
    if z_prop_parts[0].lower() not in ['d', 'p', 'n']:
        logger.warning(f"0번 필드 값이 유효하지 않음: {z_prop_parts[0]}")
        return None

    # 1번 필드 검사 (숫자 여부 확인 및 변환)
    try:
        z_first = float(z_prop_parts[1])
        z_prop_parts[1] = f"{z_first + z_base:.1f}"
    except ValueError:
        logger.warning(f"1번 필드 값이 유효하지 않음 (숫자가 아님): {z_prop_parts[1]}")
        return None

    # 2번 필드 검사 (숫자 여부 확인 및 변환)
    if len(z_prop_parts) > 2:
        try:
            z_second = float(z_prop_parts[2])
            z_prop_parts[2] = f"{z_second}"
        except ValueError:
            logger.warning(f"2번 필드 값이 유효하지 않음 (숫자가 아님): {z_prop_parts[2]}")
            return None

    # 3번 필드 검사 (텍스트 스트링, 생략 가능)
    if len(z_prop_parts) > 3:
        if not z_prop_parts[3].isalnum() and "." not in z_prop_parts[3]:
            logger.warning(f"3번 필드 값이 유효하지 않음: {z_prop_parts[3]}")
            return None

    # 최종 z_property 반환
    return ", ".join(z_prop_parts)



def save_shapes_to_txt(prs, output_file="c:\\tmp_freecad\\ppt_freecad.txt"):
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    slides = list(prs.slides)

    # 첫 슬라이드에서 scale 값 추출
    first_slide = slides[0]
    scale = extract_scale(first_slide)
    logger.info(f"첫 슬라이드에서 추출한 scale 값: {scale}")
    
    x_min, y_min = find_min_coordinates(first_slide.shapes, slide_height)

    with open(output_file, "w", encoding="utf-8") as f:
        for slide_index, slide in enumerate(slides):
            contains_freecad = any(
                shape.has_text_frame and "@freecad" in shape.text_frame.text.lower()
                for shape in slide.shapes
            )

            if not contains_freecad:
                message = f"# 슬라이드 {slide_index + 1}에 '@freecad' 없음. 종료합니다."
                logger.info(message)
                f.write(message + "\n")
                return output_file  # 반환값 추가

            z_base = extract_z_base(slide)
            header = f"# 슬라이드 {slide_index + 1} (z_base={z_base}, scale={scale})"
            logger.info(header)
            f.write(header + "\n")

            p_shapes = []
            other_shapes = []
            result_lines = []

            for shape in slide.shapes:
                try:
                    if not is_solid_line(shape):
                        message = f"# 실선이 아닌 도형 무시: {shape.name}"
                        logger.info(message)
                        f.write(message + "\n")
                        continue

                    if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and shape.auto_shape_type in [MSO_AUTO_SHAPE_TYPE.RECTANGLE, MSO_AUTO_SHAPE_TYPE.OVAL]:
                        center_x, center_y = calculate_center_coordinates(shape, slide_height, x_min, y_min)
                        rotation = get_shape_rotation(shape)
                        color = get_shape_color(shape)

                        if shape.auto_shape_type == MSO_AUTO_SHAPE_TYPE.RECTANGLE:
                            x_size = ppt_to_mm(shape.width)
                            y_size = ppt_to_mm(shape.height)
                            base_info = f"RECTANGLE,{center_x:.1f},{center_y:.1f},{x_size:.1f},{y_size:.1f},{rotation:.1f},{color}"
                        elif shape.auto_shape_type == MSO_AUTO_SHAPE_TYPE.OVAL:
                            radius = ppt_to_mm(shape.width / 2)
                            base_info = f"CIRCLE,{center_x:.1f},{center_y:.1f},{radius:.2f},{color}"
                        else:
                            continue

                        if not shape.has_text_frame or not shape.text_frame.text.strip():
                            message = "      # 경고: z_property 값이 없어서 무시합니다."
                            logger.warning(message)
                            f.write(message + "\n")
                            continue

                        z_property_original = shape.text_frame.text.strip().upper()
                        z_prop_parts = [part.strip() for part in z_property_original.split(",")]

                        z_property_original = shape.text_frame.text.strip().upper()
                        z_property = validate_and_adjust_z_property(z_property_original, z_base)
                        if z_property is None:
                            message = f"      # 경고: 유효하지 않은 z_property 값: {z_property_original}"
                            logger.warning(message)
                            f.write(message + "\n")
                            continue

                        # validate_and_adjust_z_property에서 반환된 z_property를 사용
                        result_line = f"{z_property}, {base_info}"  # z_property를 반환값으로 사용
                        if z_property.startswith("P"):
                            p_shapes.append(result_line)
                        else:
                            other_shapes.append(result_line)

                        

                except Exception as e:
                    logger.error(f"도형 처리 중 오류 발생: {e}")
                    f.write(f"# 도형 처리 중 오류 발생: {e}\n")

            # 헤더 작성
            f.write("# P/N\tz0\tz_size\tRECTANGLE\tx_center\ty_center\tx_size\ty_size\tangle\tcolor\n")
            f.write("# P/N\tz0\tz_size\tCIRCLE\tx_center\ty_center\tradius\tcolor\n")

            # 결과 정렬 및 작성
            result_lines.extend(sorted(p_shapes))
            result_lines.extend(sorted(other_shapes))
            for line in result_lines:
                line = modify_data(line, scale)
                f.write(line + "\n")

    return output_file



def main(ppt_file):
    if not os.path.exists(ppt_file) or not ppt_file.endswith(".pptx"):
        logger.error("오류: 유효한 PPTX 파일을 입력하세요.")
        return

    prs = Presentation(ppt_file)  # PPT 파일 열기
    output_file = save_shapes_to_txt(prs)  # 도형 정보를 추출하고 파일 저장

    input(f"\n>> Freecad 매크로 파일 입력 자료를 {output_file}에 저장하였습니다.")


if __name__ == "__main__":
    ppt_file = "c:\\tmp_freecad\\tmp.pptx"  # 입력 파일 경로 고정
    if not os.path.exists(ppt_file):
        logger.error(f"오류: 파일이 존재하지 않습니다: {ppt_file}")
        sys.exit(1)  # 실행 종료
    main(ppt_file)
    
