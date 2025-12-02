# 설명서
DESCRIPTION = """
코드 설명서:
1. PowerPoint 파일의 첫 번째 슬라이드만을 처리합니다.
2. 위쪽 방향으로 그려진, x축 크기가 0인 화살표만을 유효 화살표로 인식하고 그래프 곡선에 일치하도록 처리합니다.
3. 슬라이드 내 사각형의 크기를 기준으로 좌표 보정 범위를 설정하여, 해당 범위 내에서 좌표가 조정됩니다.
4. 슬라이드 하단에 아래와 같은 형식의 문자 상자가 있어야 하며, 이 상자에서 보정 범위 값을 추출합니다.
    x_min=log10(1000)
    x_max=log10(30000)
    y_min=100
    y_max=150
5. 슬라이드 최상단에 위치한 텍스트 상자의 텍스트를 그래프의 타이틀로 사용합니다.
"""

import io
import sys
import os
import re
import shutil
import tempfile
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
import matplotlib.pyplot as plt
import numpy as np
from PIL import Image
import win32clipboard
import tkinter as tk
from tkinter import messagebox
from matplotlib.backends.backend_tkagg import FigureCanvasTkAgg
from scipy.interpolate import PchipInterpolator
from math import log10

global has_log10
has_log10 = False  # log10 포함 여부 플래그 (글로벌 변수)

# 전역 변수 선언
x1_rect = y1_rect = x2_rect = y2_rect = None

def open_presentation(ppt_file):
    """프레젠테이션 파일을 임시 폴더에 복사본을 생성하여 엽니다."""
    try:
        if not os.path.exists(ppt_file):
            print(f"오류: 파일이 존재하지 않습니다. 경로를 확인하세요: {ppt_file}")
            
            return None, None

        temp_dir = tempfile.gettempdir()
        temp_file = os.path.join(temp_dir, "temp_presentation_copy.pptx")
        shutil.copy2(ppt_file, temp_file)
        
        prs = Presentation(temp_file)
        print("복사본에서 프레젠테이션 파일을 성공적으로 열었습니다.")
        
        return prs, temp_file
    except Exception as e:
        print(f"프레젠테이션 파일 열기 오류: {e}")
        
        return None, None


def extract_title_from_ppt(prs):
    """PPT에서 첫 번째 슬라이드의 첫 번째 줄을 제목으로 사용합니다."""
    slide = prs.slides[0]
    topmost_text = None
    min_y = float('inf')
    
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text:
            # y 좌표가 가장 작은 텍스트 상자를 찾기 위해 top 속성을 비교
            if shape.top < min_y:
                min_y = shape.top
                topmost_text = shape.text_frame.text.splitlines()[0]  # 첫 번째 줄만 추출
    
    return topmost_text  # y 좌표가 가장 작은 텍스트 상자의 첫 번째 줄 텍스트를 반환


def extract_rectangle_size(prs):
    """사각형 크기 추출 및 보정 전 좌표 설정"""
    global x1_rect, y1_rect, x2_rect, y2_rect
    rect_width = rect_height = None
    slide = prs.slides[0]
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and shape.auto_shape_type == MSO_AUTO_SHAPE_TYPE.RECTANGLE:
            rect_width = shape.width
            rect_height = shape.height
            left = shape.left
            top = shape.top
            x1_rect, y1_rect = left, prs.slide_height - top
            x2_rect, y2_rect = left + rect_width, prs.slide_height - (top + rect_height)
            return rect_width, rect_height
    print("오류: 사각형을 찾을 수 없습니다.")
    
    return None, None




import re

def extract_correction_range(prs):
    """보정 범위 추출 - 공백 및 특수문자 제거 후 x_min, x_max, y_min, y_max 인식"""
    global has_log10
    x_min = y_min = x_max = y_max = None
    slide = prs.slides[0]

    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text:
            # 모든 공백 문자(스페이스, 탭, 줄바꿈 등)를 제거
            text = re.sub(r'\s+', '', shape.text_frame.text)

            matches = re.findall(r'(x_min|x_max|y_min|y_max)=([\d.eE+-]+|log10\(\d+\))', text)
            
            for key, value in matches:
                if "log10" in value:
                    value = eval(value)  # log10 변환
                    has_log10 = True
                else:
                    value = float(value)

                if key == "x_min":
                    x_min = value
                elif key == "x_max":
                    x_max = value
                elif key == "y_min":
                    y_min = value
                elif key == "y_max":
                    y_max = value

            if x_min is not None and x_max is not None and y_min is not None and y_max is not None:
                if x_min >= x_max:
                    print("오류: x_min이 x_max보다 크거나 같습니다. 보정 범위를 확인하세요.")
                    return None, None, None, None
                if y_min >= y_max:
                    print("오류: y_min이 y_max보다 크거나 같습니다. 보정 범위를 확인하세요.")
                    return None, None, None, None

                return x_min, y_min, x_max, y_max

    print("오류: 보정 범위 값을 찾을 수 없습니다. 기본값을 사용합니다.")
    x_min, y_min, x_max, y_max = log10(1000), 100, log10(30000), 150
    has_log10 = True  # 기본값에 log10 포함
    return x_min, y_min, x_max, y_max




def process_line_shape(shape, slide_height, rect_width, rect_height, x_min, y_min, x_max, y_max, line_coordinates):
    """개별 선 도형을 처리하여 보정합니다."""
    global x1_rect, y1_rect, x2_rect, y2_rect
    left = shape.left
    top = shape.top
    width = shape.width
    height = shape.height
    start_x, start_y = left, slide_height - top
    end_x, end_y = left + width, slide_height - (top + height)

    if start_x == end_x:
        x_line = max(start_x, end_x)
        y_line = max(start_y, end_y)
        x0_rect = min(x1_rect, x2_rect)
        y0_rect = min(y1_rect, y2_rect)
        
        if rect_width and rect_height and x_min is not None and y_min is not None and x_max is not None and y_max is not None:
            corrected_x = (x_line - x0_rect) / rect_width * (x_max - x_min) + x_min
            corrected_y = (y_line - y0_rect) / rect_height * (y_max - y_min) + y_min
            line_coordinates.append((corrected_x, corrected_y))


def process_shapes(shapes, slide_height, rect_width, rect_height, x_min, y_min, x_max, y_max):
    """슬라이드의 도형을 처리하며 그룹 도형도 포함합니다."""
    line_coordinates = []
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            line_coordinates.extend(process_shapes(shape.shapes, slide_height, rect_width, rect_height, x_min, y_min, x_max, y_max))
        elif shape.shape_type == MSO_SHAPE_TYPE.LINE:
            process_line_shape(shape, slide_height, rect_width, rect_height, x_min, y_min, x_max, y_max, line_coordinates)
    return line_coordinates



def plot_interpolated_coordinates(line_coordinates, x_min, x_max, y_min, y_max, title=None, has_log10=False):
    """자유 형식 보간법을 사용하여 추출된 좌표 데이터를 그래프로 플롯합니다."""
    if len(line_coordinates) < 2:
        print("오류: 보간을 수행하기에 충분한 좌표가 없습니다.")
        return
    
    # 중복된 x 좌표 제거
    unique_coordinates = {}
    for coord in line_coordinates:
        if coord[0] not in unique_coordinates:
            unique_coordinates[coord[0]] = coord[1]
    
    # x 기준으로 정렬
    x_coords = sorted(unique_coordinates.keys())
    y_coords = [unique_coordinates[x] for x in x_coords]
    
    # 보간 처리
    try:
        from scipy.interpolate import make_interp_spline, interp1d
        
        spline_interpolator = make_interp_spline(x_coords, y_coords, k=3)
        x_new = np.linspace(min(x_coords), max(x_coords), num=500)
        y_new = spline_interpolator(x_new)
    except ValueError:
        interpolator = interp1d(x_coords, y_coords, kind='linear')
        x_new = np.linspace(min(x_coords), max(x_coords), num=500)
        y_new = interpolator(x_new)
    
    # 특수문자 제거 (Glyph 9 오류 방지)
    title = re.sub(r'[^\x20-\x7E]', '', title) if title else None
    
    root = tk.Tk()
    root.title("Interpolated Coordinate Graph")
    root.protocol("WM_DELETE_WINDOW", root.quit)
    
    if has_log10:
        x_new = [10 ** x for x in x_new]
        x_coords = [10 ** x for x in x_coords]
        x_min, x_max = 10 ** x_min, 10 ** x_max

    # 그래프 생성
    fig, ax = plt.subplots()
    ax.plot(x_coords, y_coords, 'ro', label='Original Data')
    ax.plot(x_new, y_new, 'b.', label='B-spline Interpolated Data', markersize=2)
    ax.set_xlim(x_min, x_max)
    ax.set_ylim(y_min, y_max)
    
    if has_log10:
        ax.set_xscale('log')
        
    ax.set_xlabel("X Coordinate")
    ax.set_ylabel("Y Coordinate")
    
    if title:
        ax.set_title(title)
    
    ax.grid()
    ax.legend()
    
    # 그래프를 Tkinter에 표시
    canvas = FigureCanvasTkAgg(fig, master=root)
    canvas.draw()
    canvas.get_tk_widget().pack()

    # 버튼 프레임 생성
    frame = tk.Frame(root)
    frame.pack(pady=10)

    # 원본 데이터 복사 기능
    def copy_original_data():
        data = "\n".join(f"{x:.3e}\t{y:.3e}" for x, y in zip(x_coords, y_coords))
        root.clipboard_clear()
        root.clipboard_append(data)
        root.update()

    # 보간된 좌표 데이터 복사 기능
    def copy_text_to_clipboard():
        data = "\n".join(f"{x:.3e}\t{y:.3e}" for x, y in zip(x_new, y_new))
        win32clipboard.OpenClipboard()
        win32clipboard.EmptyClipboard()
        win32clipboard.SetClipboardText(data, win32clipboard.CF_TEXT)
        win32clipboard.CloseClipboard()

    # 그래프 이미지 복사 기능
    def copy_image_to_clipboard():
        buf = io.BytesIO()
        fig.savefig(buf, format="png")
        buf.seek(0)
        image = Image.open(buf)
        output = io.BytesIO()
        image.convert("RGB").save(output, "BMP")
        data = output.getvalue()[14:]  # BMP 헤더 제거
        output.close()
        
        win32clipboard.OpenClipboard()
        win32clipboard.EmptyClipboard()
        win32clipboard.SetClipboardData(win32clipboard.CF_DIB, data)
        win32clipboard.CloseClipboard()

    # 버튼 추가
    tk.Button(frame, text="원본 데이터 복사", command=copy_original_data, bg="orange", fg="black").grid(row=0, column=0, padx=5, pady=5)
    tk.Button(frame, text="보간 데이터 복사", command=copy_text_to_clipboard, bg="green", fg="white").grid(row=0, column=1, padx=5, pady=5)
    tk.Button(frame, text="이미지 복사", command=copy_image_to_clipboard, bg="purple", fg="white").grid(row=0, column=2, padx=5, pady=5)
    tk.Button(frame, text="닫기", command=root.quit, bg="red", fg="white").grid(row=0, column=3, padx=5, pady=5)

    root.mainloop()


def display_correction_values(x_min, x_max, y_min, y_max):
    """보정 범위 값을 텍스트로 출력하는 함수"""
    print(f"x_min = {x_min}")
    print(f"x_max = {x_max}")
    print(f"y_min = {y_min}")
    print(f"y_max = {y_max}")
    

def main(ppt_file):
    # 코드 설명서 화면에 프린트
    print(DESCRIPTION)
    
    _, ext = os.path.splitext(ppt_file)
    if ext.lower() != ".pptx":
        print("오류: 지원되지 않는 파일 형식입니다. .pptx 파일을 사용하세요.")
        
        return
    
    prs, temp_file = open_presentation(ppt_file)
    if prs is None:
        return
    
    try:
        title = extract_title_from_ppt(prs)
        
        rect_width, rect_height = extract_rectangle_size(prs)
        if rect_width is None or rect_height is None:
            print("오류: 사각형 크기가 설정되지 않았습니다.")
            
            return
        
        x_min, y_min, x_max, y_max = extract_correction_range(prs)
        #화면에 표시
        display_correction_values(x_min, x_max, y_min, y_max)
            
        if x_min is None or y_min is None or x_max is None or y_max is None:
            print("오류: 보정 범위 값이 없습니다.")
            
            return
        
        slide = prs.slides[0]
        line_coordinates = process_shapes(slide.shapes, prs.slide_height, rect_width, rect_height, x_min, y_min, x_max, y_max)
        line_coordinates.sort(key=lambda coord: coord[0])
        
        plot_interpolated_coordinates(line_coordinates, x_min, x_max, y_min, y_max, title=title)
    
    finally:
        if temp_file and os.path.exists(temp_file):
            os.remove(temp_file)
            print("임시 파일이 삭제되었습니다.")


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("사용법: python script.py <ppt_file>")
    else:
        ppt_file = sys.argv[1]
        print(f"입력 파일: {ppt_file}")
        main(ppt_file)
