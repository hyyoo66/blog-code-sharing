# 이 코드 실행 후, 후속 실행 코드 
CALL_SUB_CODE_PATH = "sub_PPT_to_Freecad_macro_data.py"

import time
import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_AUTO_SHAPE_TYPE
import subprocess
from pptx.util import Pt
import math
import sys
from pptx.dml.color import RGBColor
import logging


# 로그 파일 초기화
LOG_FILE = "ppt_processor1.log"

# 기존 로그 파일 삭제
if os.path.exists(LOG_FILE):
    os.remove(LOG_FILE)

# 로깅 설정
logger = logging.getLogger(__name__)
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    handlers=[
#        logging.FileHandler(LOG_FILE, encoding='utf-8'),
        logging.StreamHandler(sys.stdout)
    ]
)



def calculate_rotated_coordinates(shape, group_shape=None):
    """
    도형의 좌표를 계산하는 통일된 함수
    
    Parameters:
        shape: 계산할 도형
        group_shape: 도형이 속한 그룹 (없으면 None)
        
    Returns:
        dict: 계산된 좌표 및 크기 정보
    """
    if group_shape is None:
        # 그룹에 속하지 않은 도형의 경우 원래 좌표 반환
        return {
            'left': shape.left,
            'top': shape.top,
            'width': shape.width,
            'height': shape.height,
            'rotation': shape.rotation if hasattr(shape, 'rotation') else 0
        }
    
    # 그룹의 중심점 계산
    group_center = complex(
        group_shape.left + group_shape.width / 2,
        group_shape.top + group_shape.height / 2
    )
    
    # 도형의 중심점 계산
    shape_center = complex(
        shape.left + shape.width / 2,
        shape.top + shape.height / 2
    )
    
    # 그룹의 회전각을 라디안으로 변환
    group_rotation_rad = math.radians(group_shape.rotation)
    rotation_factor = complex(
        math.cos(group_rotation_rad),
        math.sin(group_rotation_rad)
    )
    
    # 회전된 좌표 계산
    rotated_center = group_center + (shape_center - group_center) * rotation_factor
    
    # 최종 좌표 계산
    final_left = round(rotated_center.real - shape.width / 2)
    final_top = round(rotated_center.imag - shape.height / 2)
    
    return {
        'left': final_left,
        'top': final_top,
        'width': shape.width,
        'height': shape.height,
        'rotation': (shape.rotation + group_shape.rotation) % 360
    }


# get_absolute_coordinates 함수를 새로운 함수로 대체
def get_absolute_coordinates(group_shape, shape):
    """
    통일된 좌표 계산 함수를 사용하여 절대 좌표 반환
    """
    coords = calculate_rotated_coordinates(shape, group_shape)
    return {
        'left': coords['left'],
        'top': coords['top'],
        'width': coords['width'],
        'height': coords['height']
    }


def get_shape_bounds(shape):
    """도형의 경계 좌표를 반환"""
    left = Emu(shape.left)
    top = Emu(shape.top)
    width = Emu(shape.width)
    height = Emu(shape.height)
    return {
        'left': left,
        'top': top,
        'right': left + width,
        'bottom': top + height,
        'width': width,
        'height': height
    }
    
def get_group_member_shapes(group_shape, create_new=False, slide=None):
    """
    그룹 내 멤버 도형을 가져오는 함수
    - create_new: 새 도형 객체를 생성할지 여부
    - slide: 새 도형 생성 시 필요한 슬라이드 객체
    """
    if not create_new:
        # 원본 도형 객체 반환
        return list(group_shape.shapes)
    else:
        # 새로운 도형 객체 생성
        if not slide:
            raise ValueError("새 도형 생성 시 slide 객체가 필요합니다")

        grp_sp = group_shape._element
        shape_elements = [
            child for child in grp_sp
            if child.tag.endswith(('}sp', '}grpSp', '}pic'))
        ]
        shape_objects = []
        for child in shape_elements:
            shape_objects.append(slide.shapes._shape_factory(child))
        return shape_objects
 
    
def calculate_group_bounds(shapes_in_group):
    """
    그룹 내부 도형의 최외각 경계를 계산합니다.
    - 실선이 아닌 사각형이 있을 경우 해당 도형을 우선 기준으로 사용합니다.
    - 실선이 아닌 사각형이 없으면 그룹 내 모든 도형의 최외각 경계를 계산합니다.
    """
    # 실선이 아닌 사각형을 먼저 찾기
    non_solid_rectangles = [
        shape for shape in shapes_in_group
        if (shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and
            shape.auto_shape_type in [1, 2, 3] and  # 사각형 계열 도형
            get_line_style(shape.line) != "실선")  # 실선이 아닌 경우
    ]

    # 실선이 아닌 사각형이 있으면 해당 도형의 좌표를 기준으로 사용
    if non_solid_rectangles:
        reference_shape = non_solid_rectangles[0]
        x_min = reference_shape.left
        y_min = reference_shape.top
        x_max = reference_shape.left + reference_shape.width
        y_max = reference_shape.top + reference_shape.height
        print("실선이 아닌 사각형을 기준 도형으로 사용합니다.")
        return x_min, y_min, x_max, y_max

    # 실선이 아닌 사각형이 없으면 그룹 내 모든 도형의 최외각 경계 계산
    x_min = float('inf')
    y_min = float('inf')
    x_max = float('-inf')
    y_max = float('-inf')

    for shape in shapes_in_group:
        x_min = min(x_min, shape.left)
        y_min = min(y_min, shape.top)
        x_max = max(x_max, shape.left + shape.width)
        y_max = max(y_max, shape.top + shape.height)

    print("모든 도형의 최외각 경계를 기준으로 계산합니다.")
    return x_min, y_min, x_max, y_max


def calculate_group_center_complex_before_grouping(shapes):
    """
    그룹화 이전 도형들의 가상 박스 중심 좌표를 복소수로 계산합니다.

    Parameters:
        shapes (list): 그룹화 대상이 되는 도형 객체 리스트

    Returns:
        complex: 가상 박스 중심 좌표 (복소수)
    """
    min_left = float('inf')
    min_top = float('inf')
    max_right = float('-inf')
    max_bottom = float('-inf')

    for shape in shapes:
        shape_left = shape.left
        shape_top = shape.top
        shape_right = shape.left + shape.width
        shape_bottom = shape.top + shape.height

        min_left = min(min_left, shape_left)
        min_top = min(min_top, shape_top)
        max_right = max(max_right, shape_right)
        max_bottom = max(max_bottom, shape_bottom)

    group_center = complex(
        (min_left + max_right) / 2,
        (min_top + max_bottom) / 2
    )

    return group_center
      
 
def is_circle(shape):
    """
    도형이 원(Circle)인지 확인합니다.

    Args:
        shape (object): PowerPoint 도형 객체

    Returns:
        bool: 도형이 원인 경우 True, 그렇지 않으면 False
    """
    try:
        # 도형이 AutoShape인지 확인
        if shape.shape_type != MSO_SHAPE_TYPE.AUTO_SHAPE:
            return False

        # 도형이 타원(OVAL)인지 확인
        if shape.auto_shape_type != MSO_AUTO_SHAPE_TYPE.OVAL:
            return False

        # 가로와 세로 크기가 동일한 경우에만 원으로 판단
        if shape.width == shape.height:
            return True

        return False
    except AttributeError as e:
        # 도형이 예상하지 못한 속성을 가지지 않은 경우 처리
        print(f"도형 속성 확인 중 오류 발생: {e}")
        return False
'''
    # 실선이 아닌 사각형 중 텍스트가 비어있지 않은 도형 찾기
    rectangles_with_text = [
        shape for shape in non_solid_rectangles
        if shape.has_text_frame and shape.text_frame.text.strip()  # 텍스트가 비어있지 않은 경우
    ]
'''   

def ungroup_shape_without_creation(slide, group_shape):
    """
    그룹을 해제하고 새 도형을 생성하지 않음. 기존 도형을 상위 계층으로 이동.
    """
    # 그룹의 XML 요소와 부모 슬라이드의 요소 트리 가져오기
    grp_sp = group_shape._element
    sld_spTree = slide.shapes._spTree
    idx = sld_spTree.index(grp_sp)

    # 그룹 내부 도형 추출
    shape_elements = [child for child in grp_sp if child.tag.endswith(('}sp', '}grpSp', '}pic'))]

    if not shape_elements:
        print("그룹 내 도형이 없습니다. 중단합니다.")
        return []

    # 그룹 내부 도형을 슬라이드로 옮기기
    for sp in reversed(shape_elements):  # 요소 순서를 유지하면서 삽입
        sld_spTree.insert(idx, sp)

    # 그룹 도형 삭제
    sld_spTree.remove(grp_sp)

    # 기존 도형 객체를 반환 (새로 생성하지 않음)
    ungrouped_shapes = [slide.shapes._shape_factory(sp) for sp in shape_elements]
    
    print(f"그룹 '{group_shape.name}' 해제 완료. {len(ungrouped_shapes)}개의 도형이 상위 계층으로 이동했습니다.")
    return ungrouped_shapes


def get_non_solid_rectangle_info(shape):
    """
    실선이 아닌 사각형 도형 정보를 반환.
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    # 사각형인지 확인
    if shape.shape_type != MSO_SHAPE_TYPE.AUTO_SHAPE:
        return None

    # 실선 여부 확인
    if not hasattr(shape, "line") or get_line_style(shape.line) in ["실선", "없음"]:
        return None

    # 도형 정보 반환
    return {
        'name': shape.name,
        'left': shape.left,
        'top': shape.top,
        'width': shape.width,
        'height': shape.height,
        'rotation': shape.rotation,
        'text': shape.text_frame.text if shape.has_text_frame else None
    }


def calculate_bounding_box(shapes):
    """
    그룹 멤버 도형의 최외곽 경계를 계산합니다.

    Parameters:
        shapes (list): 그룹 멤버 도형 리스트

    Returns:
        tuple: (x_min, y_min, x_max, y_max)
    """
    x_min, y_min = float('inf'), float('inf')
    x_max, y_max = float('-inf'), float('-inf')

    for shape in shapes:
        # 도형의 네 꼭짓점 계산
        corners = [
            (shape.left, shape.top),
            (shape.left + shape.width, shape.top),
            (shape.left, shape.top + shape.height),
            (shape.left + shape.width, shape.top + shape.height),
        ]

        for x, y in corners:
            x_min = min(x_min, x)
            y_min = min(y_min, y)
            x_max = max(x_max, x)
            y_max = max(y_max, y)

    return x_min, y_min, x_max, y_max

def ungroup_shape(slide, group_shape):
    """
    그룹 해제 후 도형의 좌표와 회전을 복소수 연산으로 정확히 변환.
    그룹의 확대/축소를 고려하여 멤버 도형의 크기와 위치를 보정합니다.
    - 선 도형 등 모든 구성 요소가 누락되지 않도록 복원.
    
    Parameters:
        slide: 현재 슬라이드 객체
        group_shape: 해제할 그룹 도형 객체
        
    Returns:
        tuple: (ungrouped_shapes: 해제된 도형 객체 리스트, reference_shape: 기준 도형 객체)
    """
    import copy
    import math

    # 1. 그룹의 현재 절대 좌표 및 회전 각도 계산
    group_left_abs = group_shape.left
    group_top_abs = group_shape.top
    group_width_abs = group_shape.width
    group_height_abs = group_shape.height
    group_rotation_rad = math.radians(group_shape.rotation)

    # 2. 그룹 내 도형 추출 및 언그룹 (XML 수준에서 슬라이드로 이동)
    grp_sp = group_shape._element
    sld_spTree = slide.shapes._spTree
    idx = sld_spTree.index(grp_sp)

    # ✅ 모든 하위 도형 태그 포함 (선 도형 포함)
    shape_elements = [
        child for child in grp_sp
        if child.tag.endswith(('}sp', '}grpSp', '}pic', '}cxnSp'))  # 🔧 'cxnSp'는 선 도형
    ]

    if not shape_elements:
        print("그룹 내 도형이 없습니다. 그룹 해제를 중단합니다.")
        return [], None

    # 3. 그룹 해제 후 도형 상위로 이동하고, 임시 ungrouped_shapes 리스트 생성
    temp_ungrouped_shapes = []
    for sp in reversed(shape_elements):  # 역순으로 삽입하여 Z-order 유지
        sld_spTree.insert(idx, sp)
        try:
            new_shape = slide.shapes._shape_factory(sp)
            temp_ungrouped_shapes.append(new_shape)
        except Exception as e:
            print(f"⚠ 도형 생성 실패: {e}")

    # 4. 원본 그룹 도형 삭제
    sld_spTree.remove(grp_sp)
    print(f"그룹 '{group_shape.name}'의 XML 요소가 해제되었습니다.")

    # 5. 임시 ungrouped_shapes의 "가상" (unscaled) 바운딩 박스 계산
    x_min_unscaled_content, y_min_unscaled_content, \
    x_max_unscaled_content, y_max_unscaled_content = calculate_bounding_box(temp_ungrouped_shapes)

    unscaled_content_width = x_max_unscaled_content - x_min_unscaled_content
    unscaled_content_height = y_max_unscaled_content - y_min_unscaled_content
    
    print(f"Ungrouped shapes의 초기 바운딩 박스 (unscaled): "
          f"({x_min_unscaled_content}, {y_min_unscaled_content}) - "
          f"({x_max_unscaled_content}, {y_max_unscaled_content}), "
          f"폭: {unscaled_content_width}, 높이: {unscaled_content_height}")

    # 6. 스케일 팩터 계산
    scale_x = 1.0
    scale_y = 1.0

    if unscaled_content_width > 0:
        scale_x = group_width_abs / unscaled_content_width
    if unscaled_content_height > 0:
        scale_y = group_height_abs / unscaled_content_height
    
    print(f"그룹 '{group_shape.name}'의 스케일 팩터: X={scale_x:.2f}, Y={scale_y:.2f}")

    # 7. 기준 도형 찾기 (get_non_solid_rectangle_info 활용)
    reference_shape = None
    for shape in temp_ungrouped_shapes:
        # get_non_solid_rectangle_info 함수가 없을 경우를 대비한 기본 로직
        try:
            shape_info = get_non_solid_rectangle_info(shape)
            if shape_info:
                reference_shape = shape
                break
        except:
            # 함수가 없으면 첫 번째 도형을 기준으로 사용
            if reference_shape is None:
                reference_shape = shape

    # 8. 기준 도형의 "unscaled" 중심 또는 최외곽 사각형의 "unscaled" 중심
    if reference_shape:
        ref_center_cx_unscaled = complex(
            reference_shape.left + reference_shape.width / 2,
            reference_shape.top + reference_shape.height / 2
        )
        print(f"기준 도형 '{reference_shape.name}'의 unscaled 중심: {ref_center_cx_unscaled}")
    else:
        ref_center_cx_unscaled = complex(
            x_min_unscaled_content + unscaled_content_width / 2,
            y_min_unscaled_content + unscaled_content_height / 2
        )
        print(f"최외곽 사각형의 unscaled 중심: {ref_center_cx_unscaled}을 기준점으로 사용합니다.")

    # 9. 그룹의 절대 중심 (스케일 적용 후, 슬라이드 상의 실제 위치)
    group_center_abs_cx = complex(
        group_left_abs + group_width_abs / 2,
        group_top_abs + group_height_abs / 2
    )
    print(f"그룹의 절대 중심: {group_center_abs_cx}")

    # 10. 각 도형의 좌표 변환 및 회전, 스케일 적용
    final_ungrouped_shapes = []
    for shape in temp_ungrouped_shapes:
        # a. 도형의 "unscaled" 상대 중심 계산 (그룹 내부 기준)
        shape_center_cx_unscaled = complex(
            shape.left + shape.width / 2,
            shape.top + shape.height / 2
        )
        print(f"도형 '{shape.name}' 초기 (unscaled) 상대 중심: {shape_center_cx_unscaled}")

        # b. 그룹의 기준점으로부터의 "unscaled" 상대 오프셋 벡터
        offset_cx_unscaled = shape_center_cx_unscaled - ref_center_cx_unscaled
        print(f"도형 '{shape.name}' unscaled 오프셋: {offset_cx_unscaled}")

        # c. 스케일 적용 (오프셋 벡터 및 도형 자체의 크기)
        offset_cx_scaled = complex(offset_cx_unscaled.real * scale_x, offset_cx_unscaled.imag * scale_y)
        new_width = round(shape.width * scale_x)
        new_height = round(shape.height * scale_y)
        print(f"도형 '{shape.name}' scaled 오프셋: {offset_cx_scaled}, 새 크기: ({new_width}, {new_height})")

        # d. 회전 적용 (스케일된 오프셋 벡터를 그룹의 절대 중심 기준으로 회전)
        rotation_factor = complex(math.cos(group_rotation_rad), math.sin(group_rotation_rad))
        rotated_offset_cx = offset_cx_scaled * rotation_factor
        print(f"도형 '{shape.name}' 회전된 오프셋: {rotated_offset_cx}")

        # e. 최종 절대 중심 계산 (그룹 절대 중심 + 회전된 스케일 오프셋)
        final_center_abs_cx = group_center_abs_cx + rotated_offset_cx
        print(f"도형 '{shape.name}' 최종 절대 중심: {final_center_abs_cx}")

        # f. 최종 절대 좌표 (좌상단) 계산 및 적용
        shape.left = round(final_center_abs_cx.real - new_width / 2)
        shape.top = round(final_center_abs_cx.imag - new_height / 2)
        shape.width = new_width
        shape.height = new_height
        
        # 도형 자체의 회전과 그룹의 회전을 누적
        shape.rotation = (shape.rotation + group_shape.rotation) % 360 

        final_ungrouped_shapes.append(shape)

        print(f"도형 '{shape.name}' 처리 후 최종 정보: "
              f"좌표=({shape.left}, {shape.top}), "
              f"크기=({shape.width}, {shape.height}), "
              f"회전={shape.rotation}")

    print(f"그룹 '{group_shape.name}' 해제 및 모든 멤버 도형의 좌표/크기 보정 완료.")
    
    # 반환값을 첫 번째 코드 스타일로 맞춤
    return final_ungrouped_shapes, reference_shape



def calculate_group_center_before_grouping(shapes):
    """
    그룹화 이전 도형들의 가상 박스 중심 좌표를 계산합니다.

    Parameters:
        shapes (list): 그룹화 대상이 되는 도형 객체 리스트

    Returns:
        tuple: 가상 박스 중심 좌표 (X, Y)
    """
    # 초기화: 가상 박스의 최소/최대 좌표를 설정
    min_left = float('inf')
    min_top = float('inf')
    max_right = float('-inf')
    max_bottom = float('-inf')

    # 모든 도형의 좌표를 순회하며 가상 박스 경계 계산
    for shape in shapes:
        shape_left = shape.left
        shape_top = shape.top
        shape_right = shape.left + shape.width
        shape_bottom = shape.top + shape.height

        # 가상 박스의 경계 업데이트
        min_left = min(min_left, shape_left)
        min_top = min(min_top, shape_top)
        max_right = max(max_right, shape_right)
        max_bottom = max(max_bottom, shape_bottom)

    # 가상 박스 중심 좌표 계산
    center_x = (min_left + max_right) / 2  # 중심 X 좌표
    center_y = (min_top + max_bottom) / 2  # 중심 Y 좌표

    return center_x, center_y


def calculate_absolute_coordinates(group_left, group_top, shape_left, shape_top, group_rotation):
    """
    그룹 기준의 상대 좌표를 절대 좌표로 변환.
    - group_left, group_top: 그룹의 절대 좌표
    - shape_left, shape_top: 도형의 그룹 기준 상대 좌표
    - group_rotation: 그룹의 회전 각도 (degrees)
    """
    # 그룹 중심 계산 (그룹 기준 좌표)
    group_center_x = group_left
    group_center_y = group_top

    # 도형의 상대 좌표를 그룹 중심 기준으로 변환
    relative_x = shape_left - group_center_x
    relative_y = shape_top - group_center_y

    # 회전 각도를 라디안으로 변환
    rotation_radians = math.radians(group_rotation)

    # 회전 변환 공식 적용
    rotated_x = (relative_x * math.cos(rotation_radians)) - (relative_y * math.sin(rotation_radians))
    rotated_y = (relative_x * math.sin(rotation_radians)) + (relative_y * math.cos(rotation_radians))

    # 절대 좌표 계산
    corrected_x = group_center_x + rotated_x
    corrected_y = group_center_y + rotated_y

    return corrected_x, corrected_y


def correct_coordinates_with_rotation(group_left, group_top, shape_left, shape_top, group_rotation):
    """
    그룹의 회전 값을 고려하여 멤버 도형의 좌표를 보정
    - group_left, group_top: 그룹의 절대 좌표
    - shape_left, shape_top: 도형의 상대 좌표
    - group_rotation: 그룹의 회전 각도 (degrees)
    """
    # 회전 중심 계산 (그룹 중심점)
    group_center_x = group_left
    group_center_y = group_top

    # 도형 좌표를 그룹 중심을 기준으로 상대 좌표로 변환
    relative_x = shape_left - group_center_x
    relative_y = shape_top - group_center_y

    # 회전 각도를 라디안으로 변환
    rotation_radians = math.radians(group_rotation)

    # 회전 변환 공식 적용
    rotated_x = (relative_x * math.cos(rotation_radians)) - (relative_y * math.sin(rotation_radians))
    rotated_y = (relative_x * math.sin(rotation_radians)) + (relative_y * math.cos(rotation_radians))

    # 회전된 좌표를 절대 좌표로 변환
    corrected_x = group_center_x + rotated_x
    corrected_y = group_center_y + rotated_y

    return corrected_x, corrected_y



def get_line_type(line_format):
    """
    선 종류를 판별하는 함수
    """
    if line_format is None or not line_format.width:
        return "없음"
    elif line_format.dash_style is None:
        return "실선"
    elif line_format.dash_style == "sysDot":
        return "점선"
    elif line_format.dash_style == "sysDash":
        return "대쉬선"
    else:
        return "기타"

def find_shapes_with_line_type(ppt_file, target_line_type):
    """
    PowerPoint 파일에서 특정 선 종류를 가진 도형 찾기
    """
    presentation = Presentation(ppt_file)
    found_shapes = []

    for slide_index, slide in enumerate(presentation.slides):
        for shape in slide.shapes:
            if shape.line:  # 선 정보가 있는 경우
                line_type = get_line_type(shape.line)
                if line_type == target_line_type:
                    found_shapes.append({
                        "slide_index": slide_index + 1,
                        "shape_name": shape.name,
                        "line_type": line_type,
                        "position": (shape.left, shape.top),
                        "size": (shape.width, shape.height)
                    })

    return found_shapes
    
    
def get_line_style(line_format):
    """
    도형의 선 종류를 판별하는 함수
    """
    if line_format is None or not line_format.width:  # 선이 없는 경우
        return "없음"
    elif line_format.dash_style is None:  # 점선이 아닌 일반 선
        return "실선"
    else:
        return "실선 아님"  # 점선, 대쉬선 등



def apply_text_to_group_members(group_shape):
    """
    그룹 도형 내의 멤버에 텍스트를 적용하고 조건을 검사합니다.
    오류 발생 시 처리 중단하지 않고 로깅에 기록합니다.
    """
    try:
        # 그룹이 아닌 경우 처리하지 않음
        if group_shape.shape_type != MSO_SHAPE_TYPE.GROUP:
            return False, None

        # 🔶 텍스트 프레임이 있는 도형만 검사 대상으로 필터링 (선 도형 무시)
        shapes = [shape for shape in group_shape.shapes if shape.has_text_frame]

        logger.info(f"\n그룹 '{group_shape.name}' 내 도형들:")

        # 그룹 내 모든 도형 정보를 출력
        for shape in shapes:
            shape_info = f"  도형 이름: {shape.name}, 타입: {shape.shape_type}"
            if shape.has_text_frame:
                shape_info += f", 텍스트: '{shape.text_frame.text.strip()}'"
            else:
                shape_info += ", 텍스트: 없음"
            logger.debug(shape_info)

        # 텍스트가 있는 도형과 없는 도형 분리
        shapes_with_text = [
            shape for shape in shapes
            if shape.text_frame.text.strip()
        ]
        shapes_without_text = [
            shape for shape in shapes
            if not shape.text_frame.text.strip()
        ]

        # 🔶 조건 검사 - 텍스트 프레임 있는 도형만 대상으로 조건 검토
        if len(shapes_with_text) == 1 and len(shapes_with_text) + len(shapes_without_text) == len(shapes):
            source_shape = shapes_with_text[0]
            text_content = source_shape.text_frame.text.strip()

            logger.info(f"텍스트가 있는 도형: '{source_shape.name}', 텍스트: '{text_content}'")
            logger.info("텍스트 적용 결과:")

            # 🔶 텍스트 프레임 있는 도형만 적용 (선 도형 무시)
            for shape in shapes:
                shape.text_frame.text = text_content
                shape.text_frame.paragraphs[0].font.size = Pt(8)  # 글꼴 크기 설정
                font = shape.text_frame.paragraphs[0].font
                font.color.rgb = RGBColor(255, 0, 0)  # 빨간색 설정
                logger.debug(f"  도형 이름: {shape.name}, 텍스트: '{shape.text_frame.text}'")

            logger.info('그룹 ALL 조건 충족!')
            return True, text_content

        # 조건에 맞지 않을 경우
        logger.warning("그룹 ALL 조건 미충족.")
        return False, None

    except Exception as e:
        # 오류 발생 시 로그 기록
        logger.error(f"apply_text_to_group_members 함수 실행 중 오류 발생: {str(e)}")
        return False, None




def main(ppt_file):
    """
    PowerPoint 파일을 처리하여 '@freecad' 텍스트가 포함된 슬라이드의 그룹을 처리하고 결과를 저장합니다.
    그리고 후속 스크립트 'sub_PPT_to_Freecad_macro_data.py'를 자동 실행합니다.
    """
    try:
        prs = Presentation(ppt_file)
        output_dir = "c:\\tmp_freecad"
        os.makedirs(output_dir, exist_ok=True)
        output_file = os.path.join(output_dir, "tmp.pptx")

        logger.info(f"PowerPoint 파일 '{ppt_file}' 처리 시작.")

        for slide_index, slide in enumerate(prs.slides):
            logger.info(f"슬라이드 {slide_index + 1} 처리 중...")

            # '@freecad' 텍스트 확인
            contains_freecad = any(
                shape.has_text_frame and "@freecad" in shape.text_frame.text.lower()
                for shape in slide.shapes
            )

            if not contains_freecad:
                logger.info(f"슬라이드 {slide_index + 1}에 '@freecad' 없음. 건너뜀.")
                continue

            for shape in list(slide.shapes):
                if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    # 그룹에 대해 텍스트 적용 및 그룹 해체 수행
                    is_group_all, _ = apply_text_to_group_members(shape)
                    if is_group_all:
                        logger.info(f"그룹 '{shape.name}'이 '그룹 ALL'로 처리되었습니다.")
                    new_shapes = ungroup_shape(slide, shape)
                    logger.info(f"그룹 '{shape.name}' 해체 완료.")

        # 수정된 파일 저장
        while True:
            try:
                prs.save(output_file)
                logger.info(f"수정된 파일이 저장되었습니다: {output_file}")
                break
            except PermissionError:
                sys.stdout.write('\a')
                logger.warning("파일 저장 실패: 파일을 닫고 다시 시도하세요.")
                input("파일을 닫고 Enter 키를 눌러 다시 시도하세요.")
                continue

        print(f"\n>> PPT 파일을 그릅 해제하여 {output_file}에 저장하였습니다.")
        # 후속 스크립트 실행
        next_script = CALL_SUB_CODE_PATH
        input(f">> 이어서 {next_script}를 실행합니다.")

        try:
            logger.info(f"다음 스크립트를 실행합니다: {next_script}")
            os.system(f'python "{next_script}"')
        except Exception as e:
            logger.error(f"후속 스크립트 실행 중 오류 발생: {e}")

    except Exception as e:
        logger.error(f"main 함수 실행 중 오류 발생: {str(e)}")
        sys.exit(1)


if __name__ == "__main__":
   
    if len(sys.argv) != 2:
        input("사용법: python script.py <ppt 파일 경로>")
    else:
        ppt_file = sys.argv[1]
        main(ppt_file)
        

